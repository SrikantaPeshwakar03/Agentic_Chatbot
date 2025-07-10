import os
import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd

def load_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        doc = fitz.open(file_path)
        return "\n".join(page.get_text() for page in doc)

    elif ext == ".docx":
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".pptx":
        prs = pptx.Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    elif ext == ".csv":
        df = pd.read_csv(file_path)
        return df.to_string(index=False)

    elif ext in [".txt", ".md"]:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    else:
        raise ValueError("Unsupported file type: " + ext)
